import os
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.config import (
    OUTPUTS_DIR,
    DEFAULT_LOGO_PATH
)

# GetSlideZ Color Palette (Professional Light Theme)
PRIMARY_BLUE = RGBColor(30, 95, 219)     # #1E5FDB
PRIMARY_DARK = RGBColor(11, 42, 107)     # #0B2A6B
INK_TEXT     = RGBColor(11, 15, 25)      # #0B0F19
MUTED_TEXT   = RGBColor(91, 100, 114)    # #5B6472
SURFACE_BG   = RGBColor(245, 247, 250)   # #F5F7FA
BORDER_COLOR = RGBColor(229, 231, 235)   # #E5E7EB
WHITE        = RGBColor(255, 255, 255)   # #FFFFFF
ACCENT_CYAN  = RGBColor(56, 189, 248)   # #38BDF8

FONT_TITLE = "Inter"
FONT_BODY  = "Inter"

def apply_background(slide, color: RGBColor = WHITE):
    """Set background color for slide to guarantee logo and text contrast."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_text_frame_margins(tf, top=0.05, bottom=0.05, left=0.05, right=0.05):
    """Zero out padding to prevent text collision or accidental wrapping."""
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.word_wrap = True

def resolve_logo_path(custom_logo_path: Optional[str] = None) -> Optional[str]:
    """Resolve active logo path: custom uploaded logo if valid, else default logo."""
    if custom_logo_path and os.path.exists(custom_logo_path):
        return custom_logo_path
    if DEFAULT_LOGO_PATH.exists():
        return str(DEFAULT_LOGO_PATH)
    return None

def add_header_and_footer(slide, current_slide: int, total_slides: int, slide_title: str = "", company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    """
    Precision layout positioning:
    - Top-Right Logo
    - Top-Left Header Eyebrow ({COMPANY_NAME} • PRESENTATION)
    - Bottom-Left Footer ('{COMPANY_NAME} | Full Slide Title')
    - Bottom-Right Page Number
    """
    company_upper = (company_name or "ENTERPRISE SOLUTIONS").upper()

    # 1. Top-Left Eyebrow Header
    header_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.35))
    tf = header_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = f"{company_upper}  •  PRESENTATION"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = PRIMARY_BLUE

    # 2. Top-Right Logo
    active_logo = resolve_logo_path(logo_path)
    if active_logo:
        try:
            slide.shapes.add_picture(
                active_logo,
                Inches(10.8), Inches(0.3), height=Inches(0.45)
            )
        except Exception as e:
            print(f"[Logo Render Notice]: {e}")

    # 3. Bottom-Left Footer
    footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(9.5), Inches(0.4))
    tf_f = footer_box.text_frame
    set_text_frame_margins(tf_f)
    p_f = tf_f.paragraphs[0]
    full_title = slide_title.strip() if slide_title else "Executive Presentation"
    if len(full_title) > 60:
        full_title = full_title[:57] + "..."
    p_f.text = f"{company_name} |  {full_title}"
    p_f.font.size = Pt(10)
    p_f.font.name = FONT_BODY
    p_f.font.color.rgb = MUTED_TEXT

    # 4. Bottom-Right Page Number
    pg_box = slide.shapes.add_textbox(Inches(10.8), Inches(6.85), Inches(1.83), Inches(0.4))
    tf_pg = pg_box.text_frame
    set_text_frame_margins(tf_pg)
    p_pg = tf_pg.paragraphs[0]
    p_pg.alignment = PP_ALIGN.RIGHT
    p_pg.text = f"{current_slide} / {total_slides}"
    p_pg.font.size = Pt(10)
    p_pg.font.bold = True
    p_pg.font.name = FONT_BODY
    p_pg.font.color.rgb = MUTED_TEXT

def create_cover_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, SURFACE_BG)

    # Decorative Border Frame Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.6), Inches(11.93), Inches(6.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_COLOR

    # Top-Right Logo on Cover
    active_logo = resolve_logo_path(logo_path)
    if active_logo:
        try:
            slide.shapes.add_picture(active_logo, Inches(9.6), Inches(0.9), height=Inches(0.55))
        except Exception:
            pass

    # Title
    title_text = slide_data.get("title", f"{company_name} Executive Presentation")
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(8.0), Inches(2.0))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    font_size = 36
    if len(title_text) > 40:
        font_size = 28
    if len(title_text) > 70:
        font_size = 22
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = PRIMARY_DARK

    # Subtitle
    sub_text = slide_data.get("subtitle", "Executive Solution & Strategy Deck")
    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(8.0), Inches(1.3))
    tf_sub = sub_box.text_frame
    set_text_frame_margins(tf_sub)
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = sub_text
    sub_size = 18 if len(sub_text) < 70 else 14
    p_sub.font.size = Pt(sub_size)
    p_sub.font.name = FONT_BODY
    p_sub.font.color.rgb = PRIMARY_BLUE

    # Generated Hero Image Card
    img_path = slide_data.get("image_path")
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(9.2), Inches(2.0), Inches(3.0), Inches(3.0))
        except Exception:
            pass

    # Presenter / Metadata Details
    content = slide_data.get("content", [])
    if isinstance(content, list) and content:
        meta_str = "  •  ".join(content)
    else:
        meta_str = f"Presented by {company_name}"
    if len(meta_str) > 75:
        meta_str = meta_str[:72] + "..."

    pres_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(10.0), Inches(0.4))
    tf_p = pres_box.text_frame
    set_text_frame_margins(tf_p)
    p_pres = tf_p.paragraphs[0]
    p_pres.text = f"{company_name} |  {meta_str}"
    p_pres.font.size = Pt(10)
    p_pres.font.bold = True
    p_pres.font.color.rgb = MUTED_TEXT

def create_divider_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, SURFACE_BG)
    add_header_and_footer(slide, current, total, slide_data.get("title", "Section Overview"), company_name, logo_path)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(11.93), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_COLOR
    
    idx_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.0), Inches(0.5))
    tf_i = idx_box.text_frame
    set_text_frame_margins(tf_i)
    p_idx = tf_i.paragraphs[0]
    p_idx.text = f"0{current}  —  SECTION OVERVIEW"
    p_idx.font.size = Pt(13)
    p_idx.font.bold = True
    p_idx.font.color.rgb = PRIMARY_BLUE

    title_text = slide_data.get("title", "Overview")
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.5), Inches(1.5))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    font_size = 32 if len(title_text) < 45 else (26 if len(title_text) < 80 else 20)
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = PRIMARY_DARK

    sub_text = slide_data.get("subtitle", "")
    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10.5), Inches(1.8))
    tf_sub = sub_box.text_frame
    set_text_frame_margins(tf_sub)
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = sub_text
    sub_size = 16 if len(sub_text) < 100 else 13
    p_sub.font.size = Pt(sub_size)
    p_sub.font.color.rgb = INK_TEXT

def create_content_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, WHITE)
    add_header_and_footer(slide, current, total, slide_data.get("title", ""), company_name, logo_path)

    title_text = slide_data.get("title", "")
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.85), Inches(9.8), Inches(0.8))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    t_size = 28 if len(title_text) < 45 else (22 if len(title_text) < 80 else 18)
    p.font.size = Pt(t_size)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = PRIMARY_DARK

    sub_text = slide_data.get("subtitle", "")
    offset_y = 1.65
    if sub_text:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(9.8), Inches(0.55))
        tf_sub = sub_box.text_frame
        set_text_frame_margins(tf_sub)
        p_sub = tf_sub.paragraphs[0]
        if len(sub_text) > 110:
            sub_text = sub_text[:107] + "..."
        p_sub.text = sub_text
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = MUTED_TEXT
        offset_y = 2.25

    img_path = slide_data.get("image_path")
    has_image = bool(img_path and os.path.exists(img_path))
    content_width = Inches(6.5) if has_image else Inches(11.9)

    content = slide_data.get("content", [])
    if isinstance(content, str):
        content = [content]

    bullets_box = slide.shapes.add_textbox(Inches(0.7), Inches(offset_y), content_width, Inches(4.3))
    tf_b = bullets_box.text_frame
    set_text_frame_margins(tf_b)

    items = content[:5]
    total_chars = sum(len(item) for item in items)
    
    if len(items) >= 4 or total_chars > 200:
        bullet_font_size = 12.5
        paragraph_space = 6
    else:
        bullet_font_size = 14.5
        paragraph_space = 10

    for i, item in enumerate(items):
        p_b = tf_b.add_paragraph() if i > 0 else tf_b.paragraphs[0]
        clean_item = item if len(item) <= 180 else item[:177] + "..."
        p_b.text = f"•  {clean_item}"
        p_b.font.size = Pt(bullet_font_size)
        p_b.font.name = FONT_BODY
        p_b.font.color.rgb = INK_TEXT
        p_b.space_after = Pt(paragraph_space)

    if has_image:
        try:
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(offset_y), Inches(5.1), Inches(4.0))
            frame.fill.solid()
            frame.fill.fore_color.rgb = SURFACE_BG
            frame.line.color.rgb = BORDER_COLOR
            
            slide.shapes.add_picture(img_path, Inches(7.7), Inches(offset_y + 0.2), Inches(4.7), Inches(3.6))
        except Exception as e:
            print(f"[Image Render Notice]: {e}")

def create_stats_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, WHITE)
    add_header_and_footer(slide, current, total, slide_data.get("title", "Key Metrics"), company_name, logo_path)

    title_text = slide_data.get("title", "Key Metrics & Impact")
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.85), Inches(9.8), Inches(0.8))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    t_size = 28 if len(title_text) < 45 else 22
    p.font.size = Pt(t_size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK

    sub_text = slide_data.get("subtitle", "")
    if sub_text:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(9.8), Inches(0.5))
        tf_sub = sub_box.text_frame
        set_text_frame_margins(tf_sub)
        p_sub = tf_sub.paragraphs[0]
        if len(sub_text) > 110:
            sub_text = sub_text[:107] + "..."
        p_sub.text = sub_text
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = MUTED_TEXT

    content = slide_data.get("content", [])
    if isinstance(content, str):
        content = [content]

    card_width = Inches(3.7)
    card_gap = Inches(0.4)
    start_x = Inches(0.7)
    start_y = Inches(2.35)

    for i, item in enumerate(content[:3]):
        x = start_x + i * (card_width + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, Inches(4.1))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_BG
        card.line.color.rgb = BORDER_COLOR

        parts = item.split(" ", 1) if " " in item else (item, "")
        big_num = parts[0]
        desc = parts[1] if len(parts) > 1 else ""

        num_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(0.3), card_width - Inches(0.4), Inches(1.1))
        tf_num = num_box.text_frame
        set_text_frame_margins(tf_num)
        p_num = tf_num.paragraphs[0]
        p_num.text = big_num
        num_size = 38
        if len(big_num) > 5:
            num_size = 28
        if len(big_num) > 8:
            num_size = 20
        p_num.font.size = Pt(num_size)
        p_num.font.bold = True
        p_num.font.color.rgb = PRIMARY_BLUE

        desc_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(1.5), card_width - Inches(0.4), Inches(2.3))
        tf_desc = desc_box.text_frame
        set_text_frame_margins(tf_desc)
        p_desc = tf_desc.paragraphs[0]
        clean_desc = desc if len(desc) <= 120 else desc[:117] + "..."
        p_desc.text = clean_desc
        desc_size = 14 if len(clean_desc) < 60 else 12
        p_desc.font.size = Pt(desc_size)
        p_desc.font.bold = True
        p_desc.font.color.rgb = INK_TEXT

def create_cards_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, WHITE)
    add_header_and_footer(slide, current, total, slide_data.get("title", "Capabilities"), company_name, logo_path)

    title_text = slide_data.get("title", "Key Capabilities")
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.85), Inches(9.8), Inches(0.8))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    t_size = 28 if len(title_text) < 45 else 22
    p.font.size = Pt(t_size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK

    sub_text = slide_data.get("subtitle", "")
    if sub_text:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(9.8), Inches(0.5))
        tf_sub = sub_box.text_frame
        set_text_frame_margins(tf_sub)
        p_sub = tf_sub.paragraphs[0]
        if len(sub_text) > 110:
            sub_text = sub_text[:107] + "..."
        p_sub.text = sub_text
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = MUTED_TEXT

    content = slide_data.get("content", [])
    if isinstance(content, str):
        content = [content]

    card_width = Inches(3.7)
    card_gap = Inches(0.4)
    start_x = Inches(0.7)
    start_y = Inches(2.35)

    for i, item in enumerate(content[:3]):
        x = start_x + i * (card_width + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, Inches(4.1))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_BG
        card.line.color.rgb = BORDER_COLOR

        card_title_box = slide.shapes.add_textbox(x + Inches(0.25), start_y + Inches(0.35), card_width - Inches(0.5), Inches(0.6))
        tf_ct = card_title_box.text_frame
        set_text_frame_margins(tf_ct)
        p_ct = tf_ct.paragraphs[0]
        p_ct.text = f"Pilar 0{i+1}"
        p_ct.font.size = Pt(15)
        p_ct.font.bold = True
        p_ct.font.color.rgb = PRIMARY_BLUE

        card_body_box = slide.shapes.add_textbox(x + Inches(0.25), start_y + Inches(1.0), card_width - Inches(0.5), Inches(2.8))
        tf_cb = card_body_box.text_frame
        set_text_frame_margins(tf_cb)
        p_cb = tf_cb.paragraphs[0]
        clean_item = item if len(item) <= 150 else item[:147] + "..."
        p_cb.text = clean_item
        cb_size = 13.5 if len(clean_item) < 70 else (12 if len(clean_item) < 110 else 10.5)
        p_cb.font.size = Pt(cb_size)
        p_cb.font.color.rgb = INK_TEXT

def create_closing_slide(prs, slide_data: Dict[str, Any], current: int, total: int, company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, SURFACE_BG)
    add_header_and_footer(slide, current, total, slide_data.get("title", "Closing"), company_name, logo_path)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(11.93), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_COLOR

    title_text = slide_data.get("title", f"Siap Bertransformasi Bersama {company_name}?")
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(1.5))
    tf = title_box.text_frame
    set_text_frame_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    t_size = 32 if len(title_text) < 45 else 24
    p.font.size = Pt(t_size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK

    sub_text = slide_data.get("subtitle", "Jadwalkan sesi diskusi dan konsultasi lebih lanjut.")
    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.5), Inches(1.0))
    tf_sub = sub_box.text_frame
    set_text_frame_margins(tf_sub)
    p_sub = tf_sub.paragraphs[0]
    if len(sub_text) > 110:
        sub_text = sub_text[:107] + "..."
    p_sub.text = sub_text
    p_sub.font.size = Pt(16 if len(sub_text) < 70 else 13)
    p_sub.font.color.rgb = PRIMARY_BLUE

    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.7), Inches(3.2), Inches(0.75))
    btn.fill.solid()
    btn.fill.fore_color.rgb = PRIMARY_BLUE
    btn.line.color.rgb = WHITE
    tf_btn = btn.text_frame
    set_text_frame_margins(tf_btn)
    p_btn = tf_btn.paragraphs[0]
    p_btn.alignment = PP_ALIGN.CENTER
    p_btn.text = "Hubungi Tim Presenter"
    p_btn.font.size = Pt(15)
    p_btn.font.bold = True
    p_btn.font.color.rgb = WHITE

    info_box = slide.shapes.add_textbox(Inches(4.7), Inches(4.8), Inches(6.5), Inches(0.5))
    tf_info = info_box.text_frame
    set_text_frame_margins(tf_info)
    p_info = tf_info.paragraphs[0]
    p_info.text = f"{company_name} Executive Presentation"
    p_info.font.size = Pt(13.5)
    p_info.font.bold = True
    p_info.font.color.rgb = MUTED_TEXT

def convert_pptx_to_pdf_win32(pptx_path: str, pdf_path: str) -> bool:
    """Convert PPTX to PDF using MS PowerPoint COM automation on Windows."""
    try:
        import win32com.client
        try:
            import pythoncom
            pythoncom.CoInitialize()
        except ImportError:
            pass
            
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = True
        
        abs_pptx = str(Path(pptx_path).resolve())
        abs_pdf = str(Path(pdf_path).resolve())
        
        deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        deck.SaveAs(abs_pdf, 32)
        deck.Close()
        powerpoint.Quit()
        return True
    except Exception as e:
        print(f"[PowerPoint COM PDF Notice]: {e}")
        return False

def convert_pptx_to_pdf_reportlab(slides_data: List[Dict[str, Any]], pdf_path: str) -> bool:
    """Fallback PDF generator using ReportLab."""
    try:
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=landscape(letter),
            rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=20,
            textColor=colors.HexColor('#0B2A6B'),
            spaceAfter=8
        )
        sub_style = ParagraphStyle(
            'DocSubTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=13,
            textColor=colors.HexColor('#1E5FDB'),
            spaceAfter=12
        )
        body_style = ParagraphStyle(
            'DocBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=12,
            textColor=colors.HexColor('#0B0F19'),
            spaceAfter=6
        )

        story = []
        for i, s in enumerate(slides_data, 1):
            story.append(Paragraph(f"Slide {i}: {s.get('title', '')}", title_style))
            if s.get('subtitle'):
                story.append(Paragraph(f"{s.get('subtitle')}", sub_style))
            story.append(Spacer(1, 8))
            
            content = s.get('content', [])
            if isinstance(content, str):
                content = [content]
            for item in content:
                story.append(Paragraph(f"• {item}", body_style))
            story.append(Spacer(1, 18))

        doc.build(story)
        return True
    except Exception as e:
        print(f"[ReportLab PDF Notice]: {e}")
        return False

def generate_presentation_and_pdf(project_name: str, slides_data: List[Dict[str, Any]], company_name: str = "Enterprise Solutions", logo_path: Optional[str] = None) -> Dict[str, str]:
    """Generate both .pptx and .pdf files with clean alignment and custom branding."""
    clean_name = "".join(c for c in project_name if c.isalnum() or c in (' ', '_', '-')).strip()
    clean_name = clean_name.replace(' ', '_') or "slidestudioz_presentation"
    
    pptx_filename = f"{clean_name}.pptx"
    pdf_filename = f"{clean_name}.pdf"
    
    pptx_path = str(OUTPUTS_DIR / pptx_filename)
    pdf_path = str(OUTPUTS_DIR / pdf_filename)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_slides = len(slides_data)

    for i, slide_data in enumerate(slides_data, 1):
        layout = slide_data.get("layout_type", "content")
        
        if layout == "cover":
            create_cover_slide(prs, slide_data, i, total_slides, company_name, logo_path)
        elif layout == "divider":
            create_divider_slide(prs, slide_data, i, total_slides, company_name, logo_path)
        elif layout == "stats":
            create_stats_slide(prs, slide_data, i, total_slides, company_name, logo_path)
        elif layout == "cards":
            create_cards_slide(prs, slide_data, i, total_slides, company_name, logo_path)
        elif layout == "closing":
            create_closing_slide(prs, slide_data, i, total_slides, company_name, logo_path)
        else:
            create_content_slide(prs, slide_data, i, total_slides, company_name, logo_path)

    import time
    try:
        prs.save(pptx_path)
        print(f"[PPTX Generator]: Saved PPTX to {pptx_path}")
    except PermissionError:
        pptx_filename = f"{clean_name}_{int(time.time())}.pptx"
        pdf_filename = f"{clean_name}_{int(time.time())}.pdf"
        pptx_path = str(OUTPUTS_DIR / pptx_filename)
        pdf_path = str(OUTPUTS_DIR / pdf_filename)
        prs.save(pptx_path)

    pdf_success = convert_pptx_to_pdf_win32(pptx_path, pdf_path)
    if not pdf_success or not os.path.exists(pdf_path):
        convert_pptx_to_pdf_reportlab(slides_data, pdf_path)

    return {
        "pptx_path": pptx_path,
        "pdf_path": pdf_path
    }
